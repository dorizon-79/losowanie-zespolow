# --- Tryb widoku z URL + layout mobilny/ekranowy ---
import streamlit as st

view_param = (st.query_params.get("view", "") or "").lower()
locked_participant = view_param in {"ucz", "participant", "u", "p"}
locked_screen      = view_param in {"screen", "prezentacja", "ekran"}  # widok prezentacyjny

# ğŸ‘‰ EKRAN = wide, UCZESTNIK = centered, reszta = wide
layout_mode = "wide" if locked_screen else ("centered" if locked_participant else "wide")
st.set_page_config(page_title="Losowanie zespoÅ‚Ã³w", layout=layout_mode)

import pandas as pd
import random
from io import BytesIO
import unicodedata
import difflib
import qrcode
from datetime import datetime

# --- PowerPoint ---
from pptx import Presentation
from pptx.util import Inches, Pt

# =================== HasÅ‚o tylko dla organizatora ===================
ORGANIZER_PASSWORD = st.secrets.get("ORGANIZER_PASSWORD", "warsztaty")

def require_organizer_password():
    """HasÅ‚o wymagane wyÅ‚Ä…cznie w widoku organizatora."""
    if st.session_state.get("authed", False):
        return
    st.markdown("### ğŸ”’ DostÄ™p organizatora")
    with st.form("login"):
        pwd = st.text_input("HasÅ‚o", type="password", placeholder="wpisz hasÅ‚oâ€¦")
        ok = st.form_submit_button("Zaloguj")
    if ok:
        if pwd == ORGANIZER_PASSWORD:
            st.session_state["authed"] = True
            st.rerun()
        else:
            st.error("NieprawidÅ‚owe hasÅ‚o.")
    st.stop()

# KrÃ³tszy tytuÅ‚ w widoku uczestnika/ekranu, peÅ‚ny u organizatora
title_text = "ğŸ‘¥ Losowanie ZespoÅ‚Ã³w" if (locked_participant or locked_screen) else "ğŸ‘¥ Losowanie osÃ³b do zespoÅ‚Ã³w"
st.title(title_text)

# --- Usprawnienia mobilne i ekranowe ---
if locked_participant or locked_screen:
    # ğŸ‘‰ dla EKRANU: peÅ‚na szerokoÅ›Ä‡, brak max-width, responsywne fonty
    st.markdown(f"""
    <style>
      [data-testid="stToolbar"], footer {{ display: none !important; }}
      header {{ visibility: hidden; }} /* pasek tytuÅ‚u Streamlit */
      /* peÅ‚na szerokoÅ›Ä‡ kontenera */
      .block-container {{
          max-width: 100vw !important;
          padding-left: 1rem !important;
          padding-right: 1rem !important;
          padding-top: 0.5rem !important;
          padding-bottom: 2rem !important;
      }}
      [data-testid="stAppViewContainer"] {{ padding: 0 !important; }}

      /* tytuÅ‚y responsywnie */
      h1 {{ font-size: clamp(18px, 2.2vw, 36px) !important; margin: .2rem 0 1rem; }}
      .team-card h3 {{ font-size: clamp(16px, 1.4vw, 24px); margin: 0 0 .25rem 0; }}
      .team-card ul {{ margin: .25rem 0 .75rem 1.1rem; padding: 0; }}
      .team-card li {{ font-size: clamp(12px, 1.05vw, 18px); line-height: 1.35; margin: .1rem 0; }}
      @media (max-width: 768px) {{
        .block-container {{ padding-bottom: 6rem !important; }}
        h1 {{ font-size: 1.6rem !important; }}
      }}
    </style>
    """, unsafe_allow_html=True)

# Auto-scroll input na mobile (uczestnik)
if locked_participant:
    import streamlit.components.v1 as components
    components.html("""
    <script>
    window.addEventListener('load', function () {
      const root = window.parent.document;
      const inp = root.querySelector('input[type="text"]');
      if (!inp) return;
      ['focus','click'].forEach(ev => {
        inp.addEventListener(ev, () => {
          setTimeout(() => { inp.scrollIntoView({behavior:'smooth', block:'center'}); }, 150);
        });
      });
    });
    </script>
    """, height=0)

# ========= WspÃ³lny magazyn (wspÃ³Å‚dzielony miÄ™dzy sesjami) =========
@st.cache_resource
def get_store():
    return {
        "balanced_teams": None,       # list[list[dict]]
        "team_lookup": None,          # key -> {team_number, team_members}
        "all_keys": [],               # list[str]
        "display_name_map": {},       # key -> "ImiÄ™ Nazwisko" (z ogonkami)
    }
STORE = get_store()

# ======================= Pomocnicze =======================
def normalize_col(col: str) -> str:
    return col.strip().lower().replace(".", "")

def strip_accents(s: str) -> str:
    nfkd = unicodedata.normalize("NFKD", s or "")
    return "".join(c for c in nfkd if not unicodedata.combining(c))

def squash_spaces(s: str) -> str:
    return " ".join((s or "").split())

def norm_name(s: str) -> str:
    # bez ogonkÃ³w, maÅ‚e litery, zbite spacje
    return squash_spaces(strip_accents(s)).lower()

def build_keys(first_name: str, last_name: str):
    # akceptujemy "imie nazwisko" i "nazwisko imie"
    key1 = norm_name(f"{first_name} {last_name}")
    key2 = norm_name(f"{last_name} {first_name}")
    return {key1, key2}

def build_lookup_from_teams(balanced_teams):
    team_lookup, all_keys, display_name_map = {}, [], {}
    for i, team in enumerate(balanced_teams):
        for p in team:
            pretty = f"{p['ImiÄ™']} {p['Nazwisko']}".strip()
            for k in build_keys(p['ImiÄ™'], p['Nazwisko']):
                team_lookup[k] = {"team_number": i + 1, "team_members": team}
                all_keys.append(k)
                display_name_map[k] = pretty
    return team_lookup, all_keys, display_name_map

def make_qr_png(data: str) -> BytesIO:
    qr = qrcode.QRCode(version=1, box_size=8, border=2)
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf

# ---------- PowerPoint: generator prezentacji ----------
from pptx import Presentation
from pptx.util import Inches, Pt

def _font_size_for_count(n_total: int) -> int:
    if n_total <= 10:   return 36
    if n_total <= 16:   return 30
    if n_total <= 24:   return 26
    if n_total <= 32:   return 22
    return 20

def make_pptx(teams, title="Losowanie ZespoÅ‚Ã³w") -> BytesIO:
    prs = Presentation()  # 16:9
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slajd tytuÅ‚owy
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    try:
        from datetime import datetime
        slide.placeholders[1].text = datetime.now().strftime("Wyniki losowania Â· %Y-%m-%d %H:%M")
    except Exception:
        pass

    # Slajdy zespoÅ‚Ã³w
    for i, team in enumerate(teams):
        names = [f"{p['Nazwisko']} {p['ImiÄ™']}" for p in team]
        names.sort()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

        # TytuÅ‚ slajdu
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.0))
        tf = title_box.text_frame
        tf.text = f"ZespÃ³Å‚ {i+1}"
        p0 = tf.paragraphs[0]
        p0.font.size = Pt(44)
        p0.font.bold = True

        # Jedna lub dwie kolumny
        if len(names) > 14:
            split = (len(names) + 1) // 2
            columns = [names[:split], names[split:]]
            col_width = Inches(6.0)
        else:
            columns = [names]
            col_width = Inches(12.33)

        base_top = Inches(1.3)
        left = Inches(0.5)
        total = len(names)
        font_pt = _font_size_for_count(total)

        for c_idx, col_names in enumerate(columns):
            box = slide.shapes.add_textbox(left + c_idx * col_width, base_top, col_width, Inches(5.8))
            tf = box.text_frame
            tf.clear()
            for idx, nm in enumerate(col_names):
                if idx == 0:
                    p = tf.paragraphs[0]
                    p.text = nm
                else:
                    p = tf.add_paragraph()
                    p.text = nm
                p.level = 0
                p.font.size = Pt(font_pt)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ---------- Widok ekranowy (prezentacja) ----------
def _chunks(lst, n):
    for i in range(0, len(lst), n):
        yield lst[i:i+n]

def render_screen_all_teams(teams, per_row=7):
    """WyÅ›wietl wszystkie zespoÅ‚y w siatce (jedna strona) â€“ same nazwiska i imiona."""
    if not teams:
        st.warning("Brak opublikowanych wynikÃ³w.")
        return
    K = len(teams)
    per_row = max(1, min(per_row, 8 if K >= 8 else K))
    for idxs in _chunks(list(range(K)), per_row):
        cols = st.columns(len(idxs))
        for c, i in enumerate(idxs):
            with cols[c].container():
                st.markdown(f"<div class='team-card'><h3>ğŸ‘¥ ZespÃ³Å‚ {i+1}</h3></div>", unsafe_allow_html=True)
                for p in teams[i]:
                    st.markdown(f"- {p['Nazwisko']} {p['ImiÄ™']}")

def maybe_autorefresh():
    ref = st.query_params.get("refresh", "")
    try:
        sec = int(ref)
    except Exception:
        sec = 0
    if sec > 0:
        import streamlit.components.v1 as components
        components.html(f"""
        <script>
          setTimeout(function() {{ window.parent.location.reload(); }}, {sec*1000});
        </script>
        """, height=0)

def get_int_param(name: str, default: int) -> int:
    val = st.query_params.get(name, "")
    try:
        return int(val)
    except Exception:
        return default

expected_cols_map = {
    'lp': 'Lp.','nazwisko': 'Nazwisko','imiÄ™': 'ImiÄ™','imi': 'ImiÄ™',
    'stanowisko': 'Stanowisko','dziaÅ‚': 'DZIAÅ','dzial': 'DZIAÅ'
}

# =================== WybÃ³r trybu (screen tylko z URL) ===================
if locked_screen:
    mode = "ğŸ–¥ï¸ Ekran"
elif locked_participant:
    mode = "ğŸ” Uczestnik"
else:
    mode = st.radio("Wybierz tryb", ["ğŸ›ï¸ Organizator", "ğŸ” Uczestnik"])

# ========================== ORGANIZATOR ==========================
if mode == "ğŸ›ï¸ Organizator":
    require_organizer_password()

    uploaded_file = st.file_uploader("ğŸ“‚ Wybierz plik Excel (.xlsx) z listÄ… osÃ³b", type=["xlsx"])

    if uploaded_file:
        try:
            df_raw = pd.read_excel(uploaded_file)
        except Exception as e:
            st.error(f"âŒ BÅ‚Ä…d odczytu pliku: {e}")
        else:
            cleaned_cols = [normalize_col(c) for c in df_raw.columns]
            mapped_cols = { expected_cols_map[c]: df_raw.columns[i]
                            for i, c in enumerate(cleaned_cols) if c in expected_cols_map }
            required = ['Lp.', 'Nazwisko', 'ImiÄ™', 'Stanowisko', 'DZIAÅ']
            if not all(c in mapped_cols for c in required):
                st.error(f"âŒ Brakuje kolumn: {', '.join([c for c in required if c not in mapped_cols])}")
            else:
                df = df_raw.rename(columns={v:k for k,v in mapped_cols.items()})
                for col in ['ImiÄ™','Nazwisko','Stanowisko','DZIAÅ']:
                    df[col] = df[col].astype(str).map(squash_spaces)

                st.success(f"âœ… Plik wczytany. OsÃ³b: {len(df)}")
                num_teams = st.number_input("ğŸ”¢ Liczba zespoÅ‚Ã³w", 2, 20, 7)

                if st.button("ğŸ¯ Rozlosuj zespoÅ‚y"):
                    participants = df.copy()
                    N = len(participants)
                    K = num_teams

                    # Docelowe rozmiary (rÃ³Å¼nica â‰¤ 1)
                    base = N // K
                    extra = N % K
                    targets = [base + (1 if i < extra else 0) for i in range(K)]

                    teams = [[] for _ in range(K)]
                    sizes = [0] * K

                    # Tasujemy kolejnoÅ›Ä‡ dziaÅ‚Ã³w, a w kaÅ¼dym dziale â€“ osoby
                    depts = list(participants.groupby("DZIAÅ"))
                    random.shuffle(depts)

                    for dept, group in depts:
                        members = group.sample(frac=1).to_dict("records")
                        # Rundy: najpierw po jednej do zespoÅ‚Ã³w majÄ…cych wolne miejsca, potem nadwyÅ¼ki
                        while members:
                            candidates = [i for i in range(K) if sizes[i] < targets[i]]
                            if not candidates:
                                candidates = list(range(K))
                            random.shuffle(candidates)
                            for ti in candidates:
                                if not members:
                                    break
                                person = members.pop()
                                teams[ti].append(person)
                                sizes[ti] += 1

                    # Sort po nazwisku
                    for i in range(K):
                        teams[i] = sorted(teams[i], key=lambda x: x["Nazwisko"])

                    st.session_state["balanced_teams"] = teams

                # podglÄ…d + publikacja + eksporty
                if st.session_state.get("balanced_teams"):
                    teams = st.session_state["balanced_teams"]
                    K = len(teams)

                    st.markdown("### ğŸ“‹ PodglÄ…d zespoÅ‚Ã³w")
                    cols = st.columns(K)
                    for i, col in enumerate(cols):
                        col.markdown(f"### ğŸ‘¥ ZespÃ³Å‚ {i+1}")
                        for p in teams[i]:
                            col.markdown(f"- {p['Nazwisko']} {p['ImiÄ™']}")

                    if st.button("ğŸ“£ Opublikuj wyniki dla uczestnikÃ³w"):
                        lookup, keys, display_map = build_lookup_from_teams(teams)
                        STORE["balanced_teams"]   = teams
                        STORE["team_lookup"]      = lookup
                        STORE["all_keys"]         = keys
                        STORE["display_name_map"] = display_map
                        st.success("âœ… Opublikowano! PoniÅ¼ej linki i QR.")

                    if STORE["team_lookup"]:
                        st.markdown("---")
                        st.markdown("### ğŸ”— Linki i QR")

                        # Link dla uczestnikÃ³w (tylko wyszukiwarka)
                        base_url = st.text_input(
                            "Wklej adres Twojej aplikacji (bez parametrÃ³w):",
                            placeholder="https://twoja-nazwa.streamlit.app"
                        )

                        if base_url:
                            participant_url = base_url.rstrip("/") + "/?view=ucz"
                            st.markdown("**Uczestnicy (wyszukiwarka):**")
                            st.code(participant_url, language="text")
                            png_u = make_qr_png(participant_url)
                            st.image(png_u, caption="QR dla uczestnikÃ³w")
                            st.download_button("ğŸ“¥ Pobierz QR (uczestnik, PNG)", data=png_u,
                                file_name="qr_uczestnik.png", mime="image/png")

                            # Link dla ekranu/prezentacji (wszystkie zespoÅ‚y na 1 stronie)
                            cols_param = st.number_input("Ile kolumn w widoku ekranu (1â€“8)", min_value=1, max_value=8, value=min(7, K))
                            refresh_param = st.number_input("Auto-odÅ›wieÅ¼anie (sekundy, 0 = wyÅ‚Ä…czone)", min_value=0, max_value=3600, value=20)
                            screen_url = base_url.rstrip("/") + f"/?view=screen&cols={int(cols_param)}&refresh={int(refresh_param)}"
                            st.markdown("**Ekran (wszystkie zespoÅ‚y na jednej stronie):**")
                            st.code(screen_url, language="text")
                            png_s = make_qr_png(screen_url)
                            st.image(png_s, caption="QR dla ekranu/prezentacji")
                            st.download_button("ğŸ“¥ Pobierz QR (ekran, PNG)", data=png_s,
                                file_name="qr_ekran.png", mime="image/png")

                        # Eksport XLSX
                        def to_excel(teams):
                            out = BytesIO()
                            with pd.ExcelWriter(out, engine='openpyxl') as w:
                                for i, t in enumerate(teams):
                                    pd.DataFrame(t)[['Nazwisko','ImiÄ™','Stanowisko','DZIAÅ']].to_excel(
                                        w, index=False, sheet_name=f'ZespÃ³Å‚ {i+1}')
                            out.seek(0); return out
                        st.download_button("ğŸ’¾ Pobierz wyniki jako Excel",
                            to_excel(teams),
                            "wyniki_losowania.xlsx",
                            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

                        # Eksport PowerPoint
                        pptx_bytes = make_pptx(teams, title="Losowanie ZespoÅ‚Ã³w")
                        st.download_button(
                            "ğŸ“½ï¸ Pobierz prezentacjÄ™ PowerPoint (PPTX)",
                            data=pptx_bytes,
                            file_name=f"losowanie_zespolow_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )

                    if st.button("ğŸšª Wyloguj organizatora"):
                        st.session_state["authed"] = False
                        st.success("Wylogowano.")
                        st.rerun()

# ========================== EKRAN (prezentacja) ==========================
if mode == "ğŸ–¥ï¸ Ekran":
    maybe_autorefresh()
    if not STORE.get("balanced_teams"):
        st.warning("ğŸ”’ Wyniki nie sÄ… jeszcze opublikowane przez organizatora.")
    else:
        per_row = get_int_param("cols", default=7)
        render_screen_all_teams(STORE["balanced_teams"], per_row=per_row)

# ========================== UCZESTNIK ==========================
if mode == "ğŸ” Uczestnik":
    if not STORE["team_lookup"]:
        st.warning("ğŸ”’ Wyniki nie sÄ… jeszcze opublikowane przez organizatora.")
    else:
        st.subheader("ğŸ” SprawdÅº swÃ³j zespÃ³Å‚")
        full_name_in = st.text_input("Wpisz imiÄ™ i nazwisko **lub** nazwisko i imiÄ™ (dokÅ‚adnie):")
        selected_key = None
        info = None

        def norm_query(q: str) -> str:
            return norm_name(q)

        if full_name_in:
            key = norm_query(full_name_in)
            info = STORE["team_lookup"].get(key)

            if not info:
                suggestions = difflib.get_close_matches(key, STORE.get("all_keys", []), n=5, cutoff=0.75)
                if suggestions:
                    st.info("ğŸ” Nie znaleziono dokÅ‚adnego dopasowania. MoÅ¼e chodzi o:")
                    cols = st.columns(min(len(suggestions), 5))
                    for i, s in enumerate(suggestions):
                        pretty = STORE["display_name_map"].get(s, s.title())
                        if cols[i].button(pretty, key=f"sugg_{i}"):
                            selected_key = s
                else:
                    st.error("âŒ Nie znaleziono takiej osoby.")

            if selected_key:
                info = STORE["team_lookup"].get(selected_key)

        if info:
            st.success(f"âœ… JesteÅ› w Zespole {info['team_number']}")
            st.markdown("ğŸ‘¥ **SkÅ‚ad zespoÅ‚u:**")
            for m in info["team_members"]:
                st.markdown(f"- {m['Nazwisko']} {m['ImiÄ™']}")
